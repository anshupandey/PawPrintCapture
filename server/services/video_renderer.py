"""
Video rendering service for creating synchronized MP4 videos
Combines slide images with audio narration
"""

import os
import subprocess
from pathlib import Path
from typing import List, Dict, Any
import json
import tempfile

class VideoRenderer:
    def __init__(self):
        self.temp_dir = None
    
    def create_video(self, pptx_path: str, audio_files: List[Dict[str, Any]], work_dir: Path) -> str:
        """Create synchronized MP4 video from slides and audio"""
        
        try:
            # Convert slides to images
            slide_images = self._convert_slides_to_images(pptx_path, work_dir)
            
            # Create video segments for each slide
            video_segments = []
            for audio_data in audio_files:
                slide_num = audio_data['slide_number']
                audio_file = audio_data['audio_file']
                image_file = slide_images.get(slide_num)
                
                if image_file and os.path.exists(audio_file):
                    segment = self._create_video_segment(image_file, audio_file, work_dir, slide_num)
                    video_segments.append(segment)
            
            # Concatenate all segments
            final_video = self._concatenate_segments(video_segments, work_dir)
            
            return final_video
            
        except Exception as e:
            raise Exception(f"Video rendering failed: {str(e)}")
    
    def _convert_slides_to_images(self, pptx_path: str, work_dir: Path) -> Dict[int, str]:
        """Convert PowerPoint slides to high-resolution images using python-pptx and PIL"""
        
        images_dir = work_dir / "slide_images"
        images_dir.mkdir(exist_ok=True)
        
        try:
            from pptx import Presentation
            from PIL import Image, ImageDraw, ImageFont
            import io
            
            # Load the presentation
            prs = Presentation(pptx_path)
            slide_images = {}
            
            # Convert each slide to image
            for slide_idx, slide in enumerate(prs.slides):
                slide_number = slide_idx + 1
                
                # Create high-resolution image (1920x1080)
                img = Image.new('RGB', (1920, 1080), 'white')
                draw = ImageDraw.Draw(img)
                
                # Extract slide content and render as image
                try:
                    # Get slide layout information
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height
                    
                    # Scale factors for 1920x1080 output
                    scale_x = 1920 / slide_width
                    scale_y = 1080 / slide_height
                    scale = min(scale_x, scale_y)  # Maintain aspect ratio
                    
                    # Calculate positioning to center the slide
                    scaled_width = int(slide_width * scale)
                    scaled_height = int(slide_height * scale)
                    offset_x = (1920 - scaled_width) // 2
                    offset_y = (1080 - scaled_height) // 2
                    
                    # Process shapes on the slide
                    for shape in slide.shapes:
                        self._render_shape_to_image(shape, draw, offset_x, offset_y, scale)
                    
                except Exception as e:
                    print(f"Warning: Could not fully render slide {slide_number}: {e}")
                    # Create a basic slide with slide number
                    try:
                        font = ImageFont.load_default()
                        draw.text((50, 50), f"Slide {slide_number}", fill='black', font=font)
                    except:
                        draw.text((50, 50), f"Slide {slide_number}", fill='black')
                
                # Save the image
                image_path = images_dir / f"slide_{slide_number:03d}.png"
                img.save(image_path, 'PNG', dpi=(300, 300))
                slide_images[slide_number] = str(image_path)
            
            # Fallback: Use LibreOffice if python-pptx fails
            if not slide_images:
                return self._convert_slides_with_libreoffice(pptx_path, images_dir)
            
            return slide_images
            
        except ImportError:
            # Fallback to LibreOffice if python-pptx is not available
            return self._convert_slides_with_libreoffice(pptx_path, images_dir)
        except Exception as e:
            # Try LibreOffice as fallback
            try:
                return self._convert_slides_with_libreoffice(pptx_path, images_dir)
            except:
                raise Exception(f"Failed to convert slides to images: {str(e)}")
    
    def _render_shape_to_image(self, shape, draw, offset_x, offset_y, scale):
        """Render a PowerPoint shape to PIL image"""
        try:
            from pptx.shapes.base import BaseShape
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            # Get shape position and size
            left = int((shape.left * scale) + offset_x) if hasattr(shape, 'left') else 0
            top = int((shape.top * scale) + offset_y) if hasattr(shape, 'top') else 0
            width = int(shape.width * scale) if hasattr(shape, 'width') else 100
            height = int(shape.height * scale) if hasattr(shape, 'height') else 100
            
            # Handle text shapes
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text
                if text.strip():
                    try:
                        font = ImageFont.load_default()
                        draw.text((left, top), text, fill='black', font=font)
                    except:
                        draw.text((left, top), text, fill='black')
            
            # Handle other shape types (simplified)
            elif hasattr(shape, 'shape_type'):
                # Draw a placeholder rectangle for other shapes
                if width > 0 and height > 0:
                    draw.rectangle([left, top, left + width, top + height], outline='gray')
        
        except Exception as e:
            # Skip problematic shapes
            pass
    
    def _convert_slides_with_libreoffice(self, pptx_path: str, images_dir: Path) -> Dict[int, str]:
        """Fallback method using LibreOffice to convert slides to images"""
        
        try:
            # First, convert PPTX to PDF to get page count
            pdf_path = images_dir / "temp_presentation.pdf"
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(images_dir),
                pptx_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            if result.returncode != 0:
                raise Exception(f"LibreOffice PDF conversion failed: {result.stderr}")
            
            # Find the generated PDF
            pdf_files = list(images_dir.glob("*.pdf"))
            if pdf_files:
                pdf_path = pdf_files[0]
            else:
                raise Exception("PDF not created by LibreOffice")
            
            # Convert PDF pages to PNG images using ImageMagick or pdftoppm
            slide_images = {}
            try:
                # Try using pdftoppm (part of poppler-utils)
                cmd = [
                    "pdftoppm",
                    "-png",
                    "-r", "150",  # 150 DPI for good quality
                    str(pdf_path),
                    str(images_dir / "slide")
                ]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                
                if result.returncode == 0:
                    # Map generated images to slide numbers
                    for image_file in images_dir.glob("slide-*.png"):
                        try:
                            # pdftoppm creates files like "slide-1.png", "slide-2.png"
                            slide_num = int(image_file.stem.split('-')[-1])
                            slide_images[slide_num] = str(image_file)
                        except (ValueError, IndexError):
                            continue
                
            except FileNotFoundError:
                # Try ImageMagick as fallback
                try:
                    cmd = [
                        "convert",
                        "-density", "150",
                        str(pdf_path),
                        str(images_dir / "slide_%03d.png")
                    ]
                    
                    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
                    
                    if result.returncode == 0:
                        for image_file in images_dir.glob("slide_*.png"):
                            try:
                                slide_num = int(image_file.stem.split('_')[-1])
                                slide_images[slide_num] = str(image_file)
                            except (ValueError, IndexError):
                                continue
                
                except FileNotFoundError:
                    raise Exception("Neither pdftoppm nor ImageMagick found for PDF to image conversion")
            
            # Clean up PDF
            if pdf_path.exists():
                pdf_path.unlink()
            
            return slide_images
            
        except subprocess.TimeoutExpired:
            raise Exception("LibreOffice conversion timed out")
        except FileNotFoundError:
            raise Exception("LibreOffice not found. Please install LibreOffice.")
        except Exception as e:
            raise Exception(f"Failed to convert slides to images with LibreOffice: {str(e)}")
    
    def _create_video_segment(self, image_file: str, audio_file: str, work_dir: Path, slide_num: int) -> str:
        """Create a video segment from an image and audio file"""
        
        output_file = work_dir / f"segment_{slide_num:03d}.mp4"
        
        try:
            # Get audio duration
            audio_duration = self._get_audio_duration(audio_file)
            
            # Create video segment using FFmpeg
            cmd = [
                "ffmpeg",
                "-y",  # Overwrite output files
                "-loop", "1",  # Loop the image
                "-i", image_file,  # Input image
                "-i", audio_file,  # Input audio
                "-c:v", "libx264",  # Video codec
                "-tune", "stillimage",  # Optimize for still images
                "-c:a", "aac",  # Audio codec
                "-b:a", "192k",  # Audio bitrate
                "-pix_fmt", "yuv420p",  # Pixel format for compatibility
                "-shortest",  # Stop when shortest input ends
                "-t", str(audio_duration),  # Duration
                "-vf", "scale=1920:1080:force_original_aspect_ratio=decrease,pad=1920:1080:(ow-iw)/2:(oh-ih)/2:black",  # Scale and pad to 1080p
                str(output_file)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
            
            if result.returncode != 0:
                raise Exception(f"FFmpeg failed for slide {slide_num}: {result.stderr}")
            
            return str(output_file)
            
        except subprocess.TimeoutExpired:
            raise Exception(f"Video segment creation timed out for slide {slide_num}")
        except FileNotFoundError:
            raise Exception("FFmpeg not found. Please install FFmpeg.")
        except Exception as e:
            raise Exception(f"Failed to create video segment for slide {slide_num}: {str(e)}")
    
    def _get_audio_duration(self, audio_file: str) -> float:
        """Get the duration of an audio file in seconds"""
        
        try:
            cmd = [
                "ffprobe",
                "-v", "quiet",
                "-show_entries", "format=duration",
                "-of", "csv=p=0",
                audio_file
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            
            if result.returncode != 0:
                raise Exception(f"FFprobe failed: {result.stderr}")
            
            return float(result.stdout.strip())
            
        except subprocess.TimeoutExpired:
            raise Exception("Audio duration detection timed out")
        except (ValueError, FileNotFoundError) as e:
            raise Exception(f"Failed to get audio duration: {str(e)}")
    
    def _concatenate_segments(self, video_segments: List[str], work_dir: Path) -> str:
        """Concatenate video segments into final MP4"""
        
        if not video_segments:
            raise Exception("No video segments to concatenate")
        
        concat_file = work_dir / "concat_list.txt"
        output_file = work_dir / "final_video.mp4"
        
        try:
            # Create concatenation file list
            with open(concat_file, 'w') as f:
                for segment in video_segments:
                    f.write(f"file '{segment}'\n")
            
            # Concatenate using FFmpeg
            cmd = [
                "ffmpeg",
                "-y",  # Overwrite output files
                "-f", "concat",
                "-safe", "0",
                "-i", str(concat_file),
                "-c", "copy",  # Copy streams without re-encoding
                str(output_file)
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=600)
            
            if result.returncode != 0:
                raise Exception(f"Video concatenation failed: {result.stderr}")
            
            return str(output_file)
            
        except subprocess.TimeoutExpired:
            raise Exception("Video concatenation timed out")
        except FileNotFoundError:
            raise Exception("FFmpeg not found. Please install FFmpeg.")
        except Exception as e:
            raise Exception(f"Failed to concatenate video segments: {str(e)}")
    
    def add_intro_outro(self, video_file: str, intro_text: str = "", outro_text: str = "") -> str:
        """Add intro and outro slides to the video"""
        
        if not intro_text and not outro_text:
            return video_file
        
        try:
            # This would create intro/outro slides and add them to the video
            # Implementation would depend on specific requirements
            # For now, return the original video
            return video_file
            
        except Exception as e:
            print(f"Warning: Failed to add intro/outro: {e}")
            return video_file
