#!/usr/bin/env python3
"""
Main PowerPoint processing orchestrator
Coordinates all the processing steps for converting PPTX to learning modules
"""

import sys
import json
import os
import time
import traceback
import tempfile
import shutil
from pathlib import Path
from typing import Dict, Any, List
import requests

# Import from utils directory
sys.path.append(os.path.join(os.path.dirname(__file__), '..'))

from transcript_generator import TranscriptGenerator
from audio_synthesizer import AudioSynthesizer
from video_renderer import VideoRenderer
from utils.file_manager import FileManager

class PowerPointProcessor:
    def __init__(self, file_path: str, job_id: str, config: Dict[str, Any]):
        self.file_path = file_path
        self.job_id = job_id
        self.config = config
        self.work_dir = Path(tempfile.mkdtemp(prefix=f"ppt_job_{job_id}_"))
        self.file_manager = FileManager(self.work_dir)
        
        # Initialize services
        self.transcript_generator = TranscriptGenerator(config['openai_api_key'])
        self.audio_synthesizer = AudioSynthesizer(config)
        self.video_renderer = VideoRenderer()
        
        self.slides_data = []
        self.transcripts = []
        self.audio_files = []
        
    def update_job_status(self, status: str, progress: int, error_message: str = ""):
        """Update job status via API call"""
        try:
            update_data = {
                'status': status,
                'progress': progress
            }
            if error_message:
                update_data['error_message'] = error_message
            
            # Make HTTP request to update job status
            import requests
            response = requests.patch(
                f'http://localhost:5000/api/jobs/{self.job_id}',
                json=update_data,
                timeout=10
            )
            
            if response.status_code != 200:
                print(f"Failed to update job status: HTTP {response.status_code}")
                
        except Exception as e:
            print(f"Failed to update job status: {e}")

    def extract_content(self):
        """Extract text and images from PowerPoint slides with image analysis"""
        try:
            self.update_job_status('extracting', 10)
            
            # Extract slides using python-pptx and create slide images
            from pptx import Presentation
            import pytesseract
            from PIL import Image
            import io
            import base64
            
            prs = Presentation(self.file_path)
            
            # Create slide images directory
            slide_images_dir = self.work_dir / "slide_images_for_ai"
            slide_images_dir.mkdir(exist_ok=True)
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_data = {
                    'slide_number': slide_idx + 1,
                    'text_content': [],
                    'image_text': [],
                    'notes': '',
                    'slide_image_base64': None
                }
                
                # Create slide image for AI analysis
                try:
                    slide_image = self._create_slide_image(slide, slide_idx + 1)
                    if slide_image:
                        slide_image_path = slide_images_dir / f"slide_{slide_idx + 1}.png"
                        slide_image.save(slide_image_path, 'PNG', dpi=(150, 150))
                        
                        # Convert to base64 for AI analysis
                        with open(slide_image_path, 'rb') as img_file:
                            slide_data['slide_image_base64'] = base64.b64encode(img_file.read()).decode()
                except Exception as e:
                    print(f"Failed to create slide image for slide {slide_idx + 1}: {e}")
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text_content = ""
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_content += run.text
                        if text_content.strip():
                            slide_data['text_content'].append(text_content.strip())
                
                # Extract images and perform OCR
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture shape type
                        try:
                            if hasattr(shape, 'image') and hasattr(shape.image, 'blob'):
                                image = Image.open(io.BytesIO(shape.image.blob))
                                ocr_text = pytesseract.image_to_string(image).strip()
                                if ocr_text:
                                    slide_data['image_text'].append(ocr_text)
                        except Exception as e:
                            print(f"OCR failed for slide {slide_idx + 1}: {e}")
                
                # Extract notes
                if hasattr(slide, 'notes_slide') and slide.notes_slide and hasattr(slide.notes_slide, 'notes_text_frame') and slide.notes_slide.notes_text_frame:
                    slide_data['notes'] = slide.notes_slide.notes_text_frame.text.strip()
                
                self.slides_data.append(slide_data)
            
            # Convert to PDF for reference
            self.file_manager.convert_pptx_to_pdf(self.file_path)
            
            self.update_job_status('generating_transcript', 25)
            
        except Exception as e:
            error_msg = f"Content extraction failed: {str(e)}"
            self.update_job_status('error', 10, error_msg)
            raise Exception(error_msg)
    
    def _create_slide_image(self, slide, slide_number: int):
        """Create a high-quality image of a slide for AI analysis"""
        try:
            from PIL import Image, ImageDraw, ImageFont
            import io
            
            # Create a high-resolution image (1920x1080)
            img = Image.new('RGB', (1920, 1080), 'white')
            draw = ImageDraw.Draw(img)
            
            # Get slide dimensions (python-pptx uses EMUs - English Metric Units)
            slide_width = 9144000  # Standard slide width in EMUs
            slide_height = 6858000  # Standard slide height in EMUs
            
            # Calculate scale to fit slide into 1920x1080 with padding
            scale_x = 1920 / slide_width
            scale_y = 1080 / slide_height
            scale = min(scale_x, scale_y) * 0.9  # Use 90% to add some padding
            
            # Calculate positioning to center the slide
            scaled_width = int(slide_width * scale)
            scaled_height = int(slide_height * scale)
            offset_x = (1920 - scaled_width) // 2
            offset_y = (1080 - scaled_height) // 2
            
            # Draw slide background
            draw.rectangle([offset_x, offset_y, offset_x + scaled_width, offset_y + scaled_height], 
                         fill='white', outline='lightgray', width=2)
            
            # Process shapes on the slide
            for shape in slide.shapes:
                self._render_shape_to_image(shape, draw, offset_x, offset_y, scale)
            
            return img
            
        except Exception as e:
            print(f"Warning: Could not create slide image for slide {slide_number}: {e}")
            return None
    
    def _render_shape_to_image(self, shape, draw, offset_x, offset_y, scale):
        """Render a PowerPoint shape to PIL image"""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            # Get shape position and size (convert from EMUs to pixels)
            left = int((shape.left * scale) + offset_x) if hasattr(shape, 'left') else 0
            top = int((shape.top * scale) + offset_y) if hasattr(shape, 'top') else 0
            width = int(shape.width * scale) if hasattr(shape, 'width') else 100
            height = int(shape.height * scale) if hasattr(shape, 'height') else 100
            
            # Handle text shapes
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text
                if text.strip():
                    try:
                        # Try to use a better font
                        font = ImageFont.load_default()
                        # Wrap text if it's too long
                        max_width = width - 10
                        words = text.split()
                        lines = []
                        current_line = ""
                        
                        for word in words:
                            test_line = current_line + (" " if current_line else "") + word
                            bbox = draw.textbbox((0, 0), test_line, font=font)
                            if bbox[2] - bbox[0] <= max_width or not current_line:
                                current_line = test_line
                            else:
                                lines.append(current_line)
                                current_line = word
                        if current_line:
                            lines.append(current_line)
                        
                        # Draw text lines
                        for i, line in enumerate(lines):
                            draw.text((left + 5, top + 5 + i * 20), line, fill='black', font=font)
                    except Exception as font_error:
                        draw.text((left + 5, top + 5), text[:100], fill='black')
            
            # Handle other shape types (simplified visualization)
            elif hasattr(shape, 'shape_type'):
                if width > 0 and height > 0:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        # Draw a placeholder for images
                        draw.rectangle([left, top, left + width, top + height], 
                                     outline='blue', fill='lightblue', width=2)
                        draw.text((left + 5, top + 5), "[Image]", fill='darkblue')
                    else:
                        # Draw a placeholder rectangle for other shapes
                        draw.rectangle([left, top, left + width, top + height], 
                                     outline='gray', width=1)
        
        except Exception as e:
            # Skip problematic shapes
            pass

    def generate_transcripts(self):
        """Generate educational transcripts using AI"""
        try:
            self.update_job_status('generating_transcript', 30)
            
            for i, slide_data in enumerate(self.slides_data):
                transcript = self.transcript_generator.generate_slide_transcript(slide_data)
                self.transcripts.append({
                    'slide_number': slide_data['slide_number'],
                    'transcript': transcript,
                    'duration_estimate': len(transcript.split()) * 0.6  # Rough estimate: 0.6 seconds per word
                })
                
                # Update progress
                progress = 30 + (i + 1) / len(self.slides_data) * 15
                self.update_job_status('generating_transcript', int(progress))
            
            self.update_job_status('refining_transcript', 45)
            
        except Exception as e:
            error_msg = f"Transcript generation failed: {str(e)}"
            self.update_job_status('error', 30, error_msg)
            raise Exception(error_msg)

    def refine_transcripts(self):
        """Refine transcripts for better instructional design"""
        try:
            self.update_job_status('refining_transcript', 50)
            
            for i, transcript_data in enumerate(self.transcripts):
                refined_transcript = self.transcript_generator.refine_transcript(
                    transcript_data['transcript'],
                    transcript_data['slide_number']
                )
                transcript_data['transcript'] = refined_transcript
                
                # Update progress
                progress = 50 + (i + 1) / len(self.transcripts) * 10
                self.update_job_status('refining_transcript', int(progress))
            
            self.update_job_status('synthesizing_audio', 60)
            
        except Exception as e:
            error_msg = f"Transcript refinement failed: {str(e)}"
            self.update_job_status('error', 50, error_msg)
            raise Exception(error_msg)

    def synthesize_audio(self):
        """Convert transcripts to audio files"""
        try:
            self.update_job_status('synthesizing_audio', 65)
            
            for i, transcript_data in enumerate(self.transcripts):
                audio_file = self.audio_synthesizer.synthesize_text(
                    transcript_data['transcript'],
                    f"slide_{transcript_data['slide_number']}.mp3",
                    self.work_dir
                )
                
                self.audio_files.append({
                    'slide_number': transcript_data['slide_number'],
                    'audio_file': audio_file,
                    'transcript': transcript_data['transcript']
                })
                
                # Update progress
                progress = 65 + (i + 1) / len(self.transcripts) * 15
                self.update_job_status('synthesizing_audio', int(progress))
            
            self.update_job_status('embedding_audio', 80)
            
        except Exception as e:
            error_msg = f"Audio synthesis failed: {str(e)}"
            self.update_job_status('error', 65, error_msg)
            raise Exception(error_msg)

    def embed_audio_in_pptx(self):
        """Embed audio files into PowerPoint slides"""
        try:
            self.update_job_status('embedding_audio', 85)
            
            narrated_pptx = self.file_manager.embed_audio_in_slides(
                self.file_path,
                self.audio_files
            )
            
            self.update_job_status('converting_pdf', 90)
            return narrated_pptx
            
        except Exception as e:
            error_msg = f"Audio embedding failed: {str(e)}"
            self.update_job_status('error', 85, error_msg)
            raise Exception(error_msg)

    def render_video(self, narrated_pptx_path: str):
        """Create synchronized MP4 video"""
        try:
            self.update_job_status('rendering_video', 95)
            
            video_file = self.video_renderer.create_video(
                narrated_pptx_path,
                self.audio_files,
                self.work_dir
            )
            
            return video_file
            
        except Exception as e:
            error_msg = f"Video rendering failed: {str(e)}"
            self.update_job_status('error', 95, error_msg)
            raise Exception(error_msg)

    def save_outputs(self, narrated_pptx: str, video_file: str):
        """Save all output files including audio zip and update job with file paths"""
        try:
            # Create outputs directory
            outputs_dir = Path("outputs") / self.job_id
            outputs_dir.mkdir(parents=True, exist_ok=True)
            
            # Copy files to outputs directory
            final_pptx = outputs_dir / "narrated_presentation.pptx"
            final_video = outputs_dir / "learning_module.mp4"
            final_pdf = outputs_dir / "original_presentation.pdf"
            final_transcripts = outputs_dir / "transcripts.json"
            final_audio_zip = outputs_dir / "audio_files.zip"
            
            shutil.copy2(narrated_pptx, final_pptx)
            shutil.copy2(video_file, final_video)
            shutil.copy2(self.work_dir / "presentation.pdf", final_pdf)
            
            # Save transcripts as JSON
            with open(final_transcripts, 'w') as f:
                json.dump(self.transcripts, f, indent=2)
            
            # Create audio files ZIP
            self._create_audio_zip(final_audio_zip)
            
            # Update job with output file paths
            output_files = {
                'narrated_pptx': str(final_pptx),
                'video_mp4': str(final_video), 
                'pdf': str(final_pdf),
                'transcripts_json': str(final_transcripts),
                'audio_zip': str(final_audio_zip)
            }
            
            # Update job status with output files via HTTP API
            import requests
            try:
                response = requests.patch(
                    f'http://localhost:5000/api/jobs/{self.job_id}',
                    json={
                        'status': 'completed',
                        'progress': 100,
                        'output_files': output_files
                    },
                    timeout=10
                )
                if response.status_code != 200:
                    print(f"Failed to update job with output files: HTTP {response.status_code}")
            except Exception as e:
                print(f"Failed to update job with output files: {e}")
            
        except Exception as e:
            error_msg = f"Failed to save outputs: {str(e)}"
            self.update_job_status('error', 98, error_msg)
            raise Exception(error_msg)
    
    def _create_audio_zip(self, zip_path: Path):
        """Create a ZIP file containing all audio files"""
        import zipfile
        
        try:
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                for audio_data in self.audio_files:
                    audio_file = audio_data['audio_file']
                    slide_number = audio_data['slide_number']
                    
                    if os.path.exists(audio_file):
                        # Add audio file with descriptive name
                        audio_filename = f"slide_{slide_number:02d}_audio.mp3"
                        zip_file.write(audio_file, audio_filename)
                        
                        # Also create a text file with the transcript
                        transcript_filename = f"slide_{slide_number:02d}_transcript.txt"
                        transcript_content = audio_data['transcript']
                        zip_file.writestr(transcript_filename, transcript_content)
        
        except Exception as e:
            print(f"Warning: Failed to create audio ZIP: {e}")
            # Create an empty zip file to prevent errors
            with zipfile.ZipFile(zip_path, 'w') as zip_file:
                zip_file.writestr("readme.txt", "Audio files could not be packaged.")

    def cleanup(self):
        """Clean up temporary files"""
        try:
            if self.work_dir.exists():
                shutil.rmtree(self.work_dir)
        except Exception as e:
            print(f"Cleanup failed: {e}")

    def process(self):
        """Main processing pipeline"""
        try:
            self.extract_content()
            self.generate_transcripts()
            self.refine_transcripts()
            self.synthesize_audio()
            narrated_pptx = self.embed_audio_in_pptx()
            video_file = self.render_video(narrated_pptx)
            self.save_outputs(narrated_pptx, video_file)
            
            # Final status update is handled in save_outputs
            
        except Exception as e:
            print(f"Processing failed: {e}")
            traceback.print_exc()
        finally:
            self.cleanup()

def main():
    if len(sys.argv) != 4:
        print("Usage: python3 powerpoint_processor.py <file_path> <job_id> <config_json>")
        sys.exit(1)
    
    file_path = sys.argv[1]
    job_id = sys.argv[2]
    config = json.loads(sys.argv[3])
    
    processor = PowerPointProcessor(file_path, job_id, config)
    processor.process()

if __name__ == "__main__":
    main()
