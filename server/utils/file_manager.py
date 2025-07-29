"""
File management utilities for PowerPoint processing
Handles file operations, conversions, and cleanup
"""

import os
import shutil
import subprocess
from pathlib import Path
from typing import List, Dict, Any
import tempfile

class FileManager:
    def __init__(self, work_dir: Path):
        self.work_dir = work_dir
        self.work_dir.mkdir(parents=True, exist_ok=True)
    
    def convert_pptx_to_pdf(self, pptx_path: str) -> str:
        """Convert PowerPoint to PDF using LibreOffice"""
        
        output_path = self.work_dir / "presentation.pdf"
        
        try:
            cmd = [
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", str(self.work_dir),
                pptx_path
            ]
            
            result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
            
            if result.returncode != 0:
                raise Exception(f"PDF conversion failed: {result.stderr}")
            
            # LibreOffice might name the file differently, so find the PDF file
            pdf_files = list(self.work_dir.glob("*.pdf"))
            if pdf_files:
                if str(pdf_files[0]) != str(output_path):
                    shutil.move(str(pdf_files[0]), str(output_path))
                return str(output_path)
            else:
                raise Exception("PDF file not created")
                
        except subprocess.TimeoutExpired:
            raise Exception("PDF conversion timed out")
        except FileNotFoundError:
            raise Exception("LibreOffice not found. Please install LibreOffice.")
        except Exception as e:
            raise Exception(f"Failed to convert to PDF: {str(e)}")
    
    def embed_audio_in_slides(self, original_pptx: str, audio_files: List[Dict[str, Any]]) -> str:
        """Embed audio files into PowerPoint slides with auto-play functionality"""
        
        output_path = self.work_dir / "narrated_presentation.pptx"
        
        try:
            import zipfile
            import shutil
            from xml.etree import ElementTree as ET
            from pathlib import Path
            import uuid
            
            # Copy original file to output path
            shutil.copy2(original_pptx, output_path)
            
            # Create audio mapping
            audio_map = {af['slide_number']: af['audio_file'] for af in audio_files}
            
            # Work with PPTX as a ZIP file
            temp_extract = self.work_dir / "pptx_temp"
            if temp_extract.exists():
                shutil.rmtree(temp_extract)
            temp_extract.mkdir()
            
            # Extract PPTX contents
            with zipfile.ZipFile(output_path, 'r') as zip_ref:
                zip_ref.extractall(temp_extract)
            
            # Process each slide that has audio
            for slide_number, audio_file in audio_map.items():
                if os.path.exists(audio_file):
                    self._embed_audio_in_slide(temp_extract, slide_number, audio_file)
            
            # Repackage the PPTX
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
                for file_path in temp_extract.rglob('*'):
                    if file_path.is_file():
                        arc_name = file_path.relative_to(temp_extract)
                        zip_ref.write(file_path, arc_name)
            
            # Cleanup temp directory
            shutil.rmtree(temp_extract)
            
            return str(output_path)
            
        except Exception as e:
            raise Exception(f"Failed to embed audio in slides: {str(e)}")
    
    def _embed_audio_in_slide(self, pptx_dir: Path, slide_number: int, audio_file: str):
        """Embed audio file directly into a specific slide with auto-play"""
        
        try:
            slide_path = pptx_dir / "ppt" / "slides" / f"slide{slide_number}.xml"
            if not slide_path.exists():
                print(f"Warning: Slide {slide_number} not found")
                return
            
            # Copy audio file to PPTX media folder
            media_dir = pptx_dir / "ppt" / "media"
            media_dir.mkdir(exist_ok=True)
            
            audio_ext = Path(audio_file).suffix
            audio_filename = f"audio{slide_number}{audio_ext}"
            audio_dest = media_dir / audio_filename
            shutil.copy2(audio_file, audio_dest)
            
            # Parse slide XML
            tree = ET.parse(slide_path)
            root = tree.getroot()
            
            # Define namespaces
            namespaces = {
                'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
            }
            
            # Find or create cSld element
            cSld = root.find('.//p:cSld', namespaces)
            if cSld is None:
                return
            
            spTree = cSld.find('.//p:spTree', namespaces)
            if spTree is None:
                return
            
            # Create audio shape with auto-play
            audio_shape = ET.SubElement(spTree, '{http://schemas.openxmlformats.org/presentationml/2006/main}sp')
            
            # Add nvSpPr (non-visual shape properties)
            nvSpPr = ET.SubElement(audio_shape, '{http://schemas.openxmlformats.org/presentationml/2006/main}nvSpPr')
            cNvPr = ET.SubElement(nvSpPr, '{http://schemas.openxmlformats.org/presentationml/2006/main}cNvPr')
            cNvPr.set('id', str(1000 + slide_number))
            cNvPr.set('name', f'Audio {slide_number}')
            
            # Add media reference with auto-play
            audioFile = ET.SubElement(cNvPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}audioFile')
            audioFile.set('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link', f'rId{slide_number + 100}')
            
            # Add auto-play attributes
            audioFile.set('autoPlay', '1')
            audioFile.set('hideInSlideShow', '1')
            
            # Update slide relationships
            self._update_slide_relationships(pptx_dir, slide_number, audio_filename)
            
            # Save modified slide
            tree.write(slide_path, encoding='utf-8', xml_declaration=True)
            
        except Exception as e:
            print(f"Warning: Could not embed audio in slide {slide_number}: {e}")
    
    def _update_slide_relationships(self, pptx_dir: Path, slide_number: int, audio_filename: str):
        """Update slide relationships to include audio file"""
        
        try:
            rels_dir = pptx_dir / "ppt" / "slides" / "_rels"
            rels_dir.mkdir(exist_ok=True)
            
            rels_file = rels_dir / f"slide{slide_number}.xml.rels"
            
            # Create or update relationships file
            if rels_file.exists():
                tree = ET.parse(rels_file)
                root = tree.getroot()
            else:
                root = ET.Element('{http://schemas.openxmlformats.org/package/2006/relationships}Relationships')
                tree = ET.ElementTree(root)
            
            # Add audio relationship
            relationship = ET.SubElement(root, '{http://schemas.openxmlformats.org/package/2006/relationships}Relationship')
            relationship.set('Id', f'rId{slide_number + 100}')
            relationship.set('Type', 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio')
            relationship.set('Target', f'../media/{audio_filename}')
            
            # Save relationships file
            tree.write(rels_file, encoding='utf-8', xml_declaration=True)
            
        except Exception as e:
            print(f"Warning: Could not update slide relationships for slide {slide_number}: {e}")
    
    def create_scorm_package(self, output_files: Dict[str, str]) -> str:
        """Create a SCORM package with all outputs (optional feature)"""
        
        scorm_dir = self.work_dir / "scorm_package"
        scorm_dir.mkdir(exist_ok=True)
        
        try:
            # Create basic SCORM structure
            # This is a simplified implementation
            manifest_content = """<?xml version="1.0" encoding="UTF-8"?>
<manifest identifier="learning_module" version="1.2" 
          xmlns="http://www.imsproject.org/xsd/imscp_rootv1p1p2"
          xmlns:adlcp="http://www.adlnet.org/xsd/adlcp_rootv1p2">
  <metadata>
    <schema>ADL SCORM</schema>
    <schemaversion>1.2</schemaversion>
  </metadata>
  <organizations default="learning_module_org">
    <organization identifier="learning_module_org">
      <title>PowerPoint Learning Module</title>
      <item identifier="item1" identifierref="resource1">
        <title>Learning Module</title>
      </item>
    </organization>
  </organizations>
  <resources>
    <resource identifier="resource1" type="webcontent" adlcp:scormtype="sco" href="index.html">
      <file href="index.html"/>
      <file href="learning_module.mp4"/>
    </resource>
  </resources>
</manifest>"""
            
            # Write manifest
            with open(scorm_dir / "imsmanifest.xml", 'w') as f:
                f.write(manifest_content)
            
            # Create simple HTML wrapper
            html_content = """<!DOCTYPE html>
<html>
<head>
    <title>Learning Module</title>
    <meta charset="UTF-8">
</head>
<body>
    <h1>Learning Module</h1>
    <video controls width="800">
        <source src="learning_module.mp4" type="video/mp4">
        Your browser does not support the video tag.
    </video>
</body>
</html>"""
            
            with open(scorm_dir / "index.html", 'w') as f:
                f.write(html_content)
            
            # Copy video file
            if 'video_mp4' in output_files and os.path.exists(output_files['video_mp4']):
                shutil.copy2(output_files['video_mp4'], scorm_dir / "learning_module.mp4")
            
            # Create ZIP package
            scorm_zip = self.work_dir / "scorm_package.zip"
            shutil.make_archive(str(scorm_zip).replace('.zip', ''), 'zip', str(scorm_dir))
            
            return str(scorm_zip)
            
        except Exception as e:
            print(f"Warning: SCORM package creation failed: {e}")
            return ""
    
    def cleanup_temp_files(self, keep_outputs: bool = True):
        """Clean up temporary files, optionally keeping output files"""
        
        try:
            if not keep_outputs:
                shutil.rmtree(self.work_dir)
            else:
                # Remove only temporary processing files, keep outputs
                temp_patterns = [
                    "slide_images",
                    "segment_*.mp4",
                    "concat_list.txt",
                    "*.tmp"
                ]
                
                for pattern in temp_patterns:
                    for file_path in self.work_dir.glob(pattern):
                        if file_path.is_file():
                            os.remove(file_path)
                        elif file_path.is_dir():
                            shutil.rmtree(file_path)
                            
        except Exception as e:
            print(f"Warning: Cleanup failed: {e}")
    
    def validate_file_types(self, file_path: str) -> bool:
        """Validate that uploaded file is a valid PPTX"""
        
        try:
            from pptx import Presentation
            
            # Try to open the file with python-pptx
            prs = Presentation(file_path)
            
            # Basic validation - check if it has slides
            if len(prs.slides) == 0:
                return False
                
            return True
            
        except Exception:
            return False
    
    def get_file_info(self, file_path: str) -> Dict[str, Any]:
        """Get information about the uploaded PowerPoint file"""
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            return {
                'slide_count': len(prs.slides),
                'file_size': os.path.getsize(file_path),
                'has_images': any(
                    any(shape.shape_type == 13 for shape in slide.shapes)
                    for slide in prs.slides
                ),
                'has_notes': any(
                    slide.notes_slide.notes_text_frame.text.strip()
                    for slide in prs.slides
                    if slide.notes_slide.notes_text_frame
                )
            }
            
        except Exception as e:
            return {'error': str(e)}
